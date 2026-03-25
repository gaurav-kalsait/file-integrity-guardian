"""
FILE INTEGRITY GUARDIAN — PowerPoint Generator
================================================
Run this script to auto-generate a .pptx file
with dark theme, diagrams, and screenshot placeholders.

Usage:
    pip install python-pptx
    python generate_ppt.py

Output: File_Integrity_Guardian_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ━━━ COLORS (White Theme) ━━━
BG_DARK = RGBColor(0xFF, 0xFF, 0xFF)       # White background
BG_CARD = RGBColor(0xF4, 0xF6, 0xF9)       # Light gray cards
BG_INPUT = RGBColor(0xEE, 0xF0, 0xF5)      # Slightly darker gray
CYAN = RGBColor(0x00, 0x7A, 0xCC)          # Darker cyan for readability on white
GREEN = RGBColor(0x00, 0x8A, 0x4E)         # Dark green
RED = RGBColor(0xCC, 0x22, 0x44)           # Dark red
AMBER = RGBColor(0xBB, 0x77, 0x00)         # Dark amber
PURPLE = RGBColor(0x6B, 0x4F, 0xBB)        # Dark purple
WHITE = RGBColor(0x1A, 0x1A, 0x2E)         # Dark text (main)
GRAY = RGBColor(0x55, 0x5E, 0x70)          # Medium gray text
DIM = RGBColor(0x88, 0x90, 0xA0)           # Light gray text
BORDER = RGBColor(0xDD, 0xDF, 0xE6)        # Light borders

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_box(slide, left, top, width, height, fill_color=BG_CARD, border_color=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    shape.rotation = 0
    return shape

def add_screenshot_placeholder(slide, left, top, width, height, label="Paste Screenshot Here"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF9)
    shape.line.color.rgb = RGBColor(0xCC, 0xCE, 0xD6)
    shape.line.width = Pt(2)
    shape.line.dash_style = 2
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"📸 {label}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x99, 0x9E, 0xAA)
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(40)

def add_slide_title(slide, title):
    add_text(slide, 0.8, 0.3, 11, 0.6, title, font_size=28, color=CYAN, bold=True)
    # Underline
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER
    shape.line.fill.background()

def add_bullet_list(slide, left, top, width, items, color=WHITE, font_size=14):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(len(items) * 0.45))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)


# ═══════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0, 1.5, 13.3, 1, "🛡️", font_size=60, alignment=PP_ALIGN.CENTER)
add_text(slide, 0, 2.5, 13.3, 0.8, "File Integrity Guardian", font_size=44, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 0, 3.3, 13.3, 0.6, "Full-Stack File Integrity Monitoring Tool", font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, 0, 4.2, 13.3, 0.4, "SHA-256 & MD5 Cryptographic Hash Verification  •  Python + React", font_size=16, color=DIM, alignment=PP_ALIGN.CENTER)
add_text(slide, 0, 5.5, 13.3, 0.4, "Presented by: Gaurav Kumar", font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(slide, 0, 6.0, 13.3, 0.4, "Course: Business Impact Analysis / Cybersecurity", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 2: AGENDA
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "Agenda")
items_left = ["01  Introduction & Abstract", "02  Problem Statement", "03  What is FIM?", "04  Why FIM Matters",
              "05  Research — Industry Tools", "06  Technology Stack", "07  Hashing Algorithms",
              "08  SHA-256 vs MD5", "09  Impact Analysis", "10  Business Impact", "11  System Architecture"]
items_right = ["12  Data Flow Diagram", "13  Python Backend", "14  React Frontend", "15  API Endpoints",
               "16-22  Tool Demo + Screenshots", "23-25  Security Recommendations", "26  Proof of Concept",
               "27  Tools & Technologies", "28  Challenges & Learnings", "29  Future Enhancements",
               "30  References", "31  Thank You / Q&A"]
for i, item in enumerate(items_left):
    add_text(slide, 1, 1.3 + i * 0.42, 5, 0.4, item, font_size=14, color=WHITE if i % 2 == 0 else GRAY)
for i, item in enumerate(items_right):
    add_text(slide, 7, 1.3 + i * 0.42, 5, 0.4, item, font_size=14, color=WHITE if i % 2 == 0 else GRAY)


# ═══════════════════════════════════════════
# SLIDE 3: ABSTRACT
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "01 — Abstract")
for i, (label, text, color) in enumerate([
    ("GOAL", "Build a full-stack File Integrity Monitoring tool that verifies system file integrity using cryptographic hashing (SHA-256 & MD5).", CYAN),
    ("METHOD", "Python backend uses hashlib to compute real hashes on OS files, stores a trusted baseline, and compares on each scan to detect unauthorized changes.", GREEN),
    ("OUTPUT", "Interactive React dashboard with real-time scanning, file browser, auto-scheduler, 5-level status indicator, alert system, and security report generation.", AMBER),
]):
    add_box(slide, 0.8, 1.3 + i * 1.7, 11.5, 1.4)
    add_text(slide, 1.1, 1.4 + i * 1.7, 2, 0.4, label, font_size=12, color=color, bold=True)
    add_text(slide, 1.1, 1.8 + i * 1.7, 11, 0.8, text, font_size=16, color=WHITE)


# ═══════════════════════════════════════════
# SLIDE 4: PROBLEM STATEMENT
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "02 — Problem Statement")
add_box(slide, 0.8, 1.3, 11.5, 1.2, fill_color=RGBColor(0xFF, 0xF0, 0xF0), border_color=RGBColor(0xCC, 0x22, 0x44))
add_text(slide, 1.1, 1.35, 2, 0.3, "THE THREAT", font_size=11, color=RED, bold=True)
add_text(slide, 1.1, 1.7, 11, 0.6, "Attackers modify critical system files to install backdoors, escalate privileges, disable firewalls, and establish persistence — often undetected for months.", font_size=15, color=WHITE)
add_text(slide, 1.1, 2.8, 3, 0.4, "KEY CHALLENGES", font_size=11, color=DIM, bold=True)
challenges = ["No trusted baseline to compare current file states against",
              "Manual file checking is impractical at scale",
              "Unauthorized modifications to OS binaries go undetected",
              "Compliance standards (PCI DSS, HIPAA, SOX) require FIM"]
add_bullet_list(slide, 1.1, 3.2, 11, challenges, color=GRAY, font_size=14)
add_box(slide, 0.8, 5.2, 11.5, 1.2, fill_color=RGBColor(0xF0, 0xFF, 0xF0), border_color=RGBColor(0x00, 0x8A, 0x4E))
add_text(slide, 1.1, 5.3, 2, 0.3, "OUR SOLUTION", font_size=11, color=GREEN, bold=True)
add_text(slide, 1.1, 5.65, 11, 0.5, "File Integrity Guardian — automated hashing, monitoring, alerting, and reporting for critical system files.", font_size=15, color=WHITE)


# ═══════════════════════════════════════════
# SLIDE 5: WHAT IS FIM + PROCESS DIAGRAM
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "03 — What is File Integrity Monitoring?")
add_text(slide, 0.8, 1.3, 11.5, 0.8, "FIM is a cybersecurity practice that detects changes to critical system files by comparing their current state against a trusted baseline using cryptographic hash functions.", font_size=15, color=GRAY)
# Process diagram boxes
steps = [("Read File", CYAN), ("Compute Hash", CYAN), ("Store Baseline", GREEN), ("Re-scan", AMBER), ("Compare & Alert", RED)]
for i, (label, color) in enumerate(steps):
    x = 0.8 + i * 2.4
    add_box(slide, x, 2.5, 2.1, 0.8, border_color=color)
    add_text(slide, x, 2.6, 2.1, 0.6, label, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_text(slide, x + 2.1, 2.65, 0.3, 0.5, "→", font_size=20, color=DIM, alignment=PP_ALIGN.CENTER)
# Info cards
for i, (title, desc) in enumerate([
    ("What it monitors", "OS binaries, config files, logs, registries, permissions"),
    ("Standards requiring FIM", "PCI DSS (11.5), HIPAA, SOX, ISO 17799, CIS Benchmarks"),
    ("Detection types", "Content changes, permission changes, file deletion, new files"),
]):
    add_box(slide, 0.8 + i * 4, 3.8, 3.7, 1.2)
    add_text(slide, 1 + i * 4, 3.9, 3.4, 0.3, title, font_size=11, color=CYAN, bold=True)
    add_text(slide, 1 + i * 4, 4.3, 3.4, 0.6, desc, font_size=12, color=GRAY)


# ═══════════════════════════════════════════
# SLIDE 6: WHY FIM MATTERS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "04 — Why FIM Matters")
for i, (icon, title, desc, color) in enumerate([
    ("🔓", "Backdoor Detection", "Modified binaries reveal attacker persistence", RED),
    ("🛡️", "Compliance", "PCI DSS, HIPAA, SOX all mandate FIM controls", CYAN),
    ("⚡", "Early Warning", "Detect breaches in minutes, not months", AMBER),
]):
    add_box(slide, 0.8 + i * 4, 1.3, 3.7, 2.5)
    add_text(slide, 0.8 + i * 4, 1.5, 3.7, 0.6, icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text(slide, 0.8 + i * 4, 2.1, 3.7, 0.4, title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, 1 + i * 4, 2.6, 3.3, 0.8, desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
add_box(slide, 0.8, 4.2, 11.5, 1, fill_color=RGBColor(0xFF, 0xF0, 0xF0), border_color=RGBColor(0xCC, 0x22, 0x44))
add_text(slide, 1.1, 4.35, 11, 0.7, "Key Stat: 9.8% of organizations fail PCI audits due to lacking FIM (Verizon). Average breach detection without FIM: 197 days (IBM).", font_size=14, color=WHITE)


# ═══════════════════════════════════════════
# SLIDE 7: RESEARCH — INDUSTRY TOOLS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "05 — Research: Industry FIM Tools")
tools_data = [
    ("Tool", "URL", "Category", "Popularity"),
    ("OSSEC", "ossec.net", "Open-source HIDS+FIM", "10k+ GitHub stars"),
    ("Tripwire", "tripwire.com", "Commercial FIM", "Fortune 500 standard"),
    ("Wazuh", "wazuh.com", "Open-source SIEM+FIM", "30k+ GitHub stars"),
    ("SolarWinds SEM", "solarwinds.com", "Commercial SIEM+FIM", "Gartner Leader"),
    ("AIDE", "aide.github.io", "Open-source FIM", "Linux standard"),
    ("Qualys FIM", "qualys.com", "Cloud FIM", "19k+ enterprises"),
    ("CrowdStrike", "crowdstrike.com", "Cloud SaaS FIM", "EPP market leader"),
    ("Samhain", "la-samhna.de", "Open-source HIDS", "Linux/Unix standard"),
]
for i, row in enumerate(tools_data):
    y = 1.3 + i * 0.6
    is_header = i == 0
    for j, cell in enumerate(row):
        x = 0.8 + j * 3
        w = 2.8
        c = CYAN if is_header else (CYAN if j == 0 else GREEN if j == 1 else WHITE if j == 2 else GRAY)
        add_text(slide, x, y, w, 0.4, cell, font_size=12 if not is_header else 11, color=c, bold=is_header)
    if i > 0:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y + 0.5), Inches(11.5), Inches(0.01))
        shape.fill.solid(); shape.fill.fore_color.rgb = BORDER; shape.line.fill.background()


# ═══════════════════════════════════════════
# SLIDE 8: TECHNOLOGY STACK
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "06 — Technology Stack Comparison")
stacks = [
    ("OSSEC", "C, SQLite, MD5/SHA-1/SHA-256, inotify, manager-agent", CYAN),
    ("Tripwire", "C/C++, proprietary DB, SHA-256/SHA-512/HAVAL, policy-driven", PURPLE),
    ("Wazuh", "OSSEC fork, Elasticsearch/Kibana, RESTful API, real-time + scheduled", AMBER),
    ("AIDE", "C, flat-file DB, MD5/SHA-1/SHA-256/SHA-512, cron-based", GREEN),
]
for i, (name, stack, color) in enumerate(stacks):
    add_box(slide, 0.8, 1.3 + i * 1.1, 11.5, 0.85)
    add_text(slide, 1.1, 1.35 + i * 1.1, 2.5, 0.4, name, font_size=15, color=color, bold=True)
    add_text(slide, 3.5, 1.4 + i * 1.1, 8.5, 0.6, stack, font_size=12, color=GRAY)
add_box(slide, 0.8, 5.8, 11.5, 1.2, fill_color=RGBColor(0xF0, 0xFF, 0xF0), border_color=RGBColor(0x00, 0x8A, 0x4E))
add_text(slide, 1.1, 5.85, 4, 0.3, "OUR STACK", font_size=11, color=GREEN, bold=True)
add_text(slide, 1.1, 6.2, 11, 0.6, "Python 3.10+ (hashlib, Flask) + React.js (Vite) + JSON persistence + REST API + Cross-platform Windows & Linux", font_size=14, color=WHITE)


# ═══════════════════════════════════════════
# SLIDE 9: HASHING ALGORITHMS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "07 — Hashing Algorithms")
algo_data = [
    ("Algorithm", "Output Size", "Security Level", "FIM Usage"),
    ("MD5", "128-bit (32 hex)", "❌ Broken — collisions feasible", "Legacy only"),
    ("SHA-1", "160-bit (40 hex)", "⚠️ Weakened — SHAttered 2017", "Being phased out"),
    ("SHA-256", "256-bit (64 hex)", "✅ Secure — no known attacks", "Industry standard"),
    ("SHA-512", "512-bit (128 hex)", "✅ Strongest common option", "High-security"),
]
for i, row in enumerate(algo_data):
    y = 1.3 + i * 0.65
    is_header = i == 0
    for j, cell in enumerate(row):
        x = 0.8 + j * 3
        c = CYAN if is_header else (CYAN if j == 0 else WHITE)
        add_text(slide, x, y, 2.8, 0.4, cell, font_size=12 if not is_header else 11, color=c, bold=is_header)
    if i > 0:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y + 0.55), Inches(11.5), Inches(0.01))
        shape.fill.solid(); shape.fill.fore_color.rgb = BORDER; shape.line.fill.background()
# Hash process diagram
steps_h = [("Input File\n(any size)", CYAN), ("SHA-256\nhashlib.sha256()", GREEN), ("64-char hex\nfingerprint", AMBER)]
for i, (label, color) in enumerate(steps_h):
    x = 1.5 + i * 3.8
    add_box(slide, x, 5, 3.2, 1.2, border_color=color)
    add_text(slide, x, 5.1, 3.2, 1, label, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    if i < 2:
        add_text(slide, x + 3.2, 5.3, 0.6, 0.6, "→", font_size=22, color=DIM, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 10: SHA-256 vs MD5
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "08 — SHA-256 vs MD5 Deep Dive")
# MD5 box
add_box(slide, 0.8, 1.3, 5.5, 5.2, fill_color=RGBColor(0xFF, 0xF0, 0xF0), border_color=RGBColor(0xCC, 0x22, 0x44))
add_text(slide, 1.2, 1.4, 4.8, 0.5, "MD5 (Broken)", font_size=20, color=RED, bold=True)
md5_items = ["128-bit output (32 hex characters)", "Created by Ron Rivest in 1991", "Collision attacks proven feasible", "Fast but insecure for FIM", "Used only for backward compatibility", "NOT suitable for security use"]
for i, item in enumerate(md5_items):
    add_text(slide, 1.2, 2.1 + i * 0.6, 4.8, 0.4, f"✗  {item}", font_size=13, color=WHITE)
# SHA box
add_box(slide, 6.8, 1.3, 5.5, 5.2, fill_color=RGBColor(0xF0, 0xFF, 0xF0), border_color=RGBColor(0x00, 0x8A, 0x4E))
add_text(slide, 7.2, 1.4, 4.8, 0.5, "SHA-256 (Secure)", font_size=20, color=GREEN, bold=True)
sha_items = ["256-bit output (64 hex characters)", "Designed by NSA in 2001 (SHA-2)", "No known practical collisions", "Used in Bitcoin, SSL/TLS, code signing", "Industry standard for modern FIM", "Our primary algorithm ✓"]
for i, item in enumerate(sha_items):
    add_text(slide, 7.2, 2.1 + i * 0.6, 4.8, 0.4, f"✓  {item}", font_size=13, color=WHITE)


# ═══════════════════════════════════════════
# SLIDE 11: IMPACT ANALYSIS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "09 — Impact Analysis: Attack Targets")
targets = [
    ("/etc/passwd & shadow", "Create backdoor accounts, weaken passwords", "CRITICAL", RED),
    ("cmd.exe / powershell.exe", "Replace with trojanized versions capturing creds", "CRITICAL", RED),
    ("svchost.exe / kernel32.dll", "DLL hijacking, rootkit installation", "CRITICAL", RED),
    ("hosts file", "DNS poisoning, redirect traffic to malicious sites", "HIGH", AMBER),
    ("crontab / scheduled tasks", "Establish persistence via scheduled malicious tasks", "HIGH", AMBER),
    ("firewall rules (iptables)", "Disable network protection, open access", "CRITICAL", RED),
    ("SSH config / sshd_config", "Weaken SSH security, allow unauthorized access", "HIGH", AMBER),
    ("/boot/vmlinuz (kernel)", "Install kernel-level rootkits", "CRITICAL", RED),
]
for i, (file, impact, risk, color) in enumerate(targets):
    y = 1.3 + i * 0.7
    add_text(slide, 0.8, y, 3.2, 0.4, file, font_size=12, color=CYAN, bold=True, font_name="Consolas")
    add_text(slide, 4.2, y, 6.5, 0.4, impact, font_size=12, color=GRAY)
    add_text(slide, 11, y, 1.5, 0.4, risk, font_size=10, color=color, bold=True)
    if i < 7:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y + 0.6), Inches(11.8), Inches(0.01))
        shape.fill.solid(); shape.fill.fore_color.rgb = BORDER; shape.line.fill.background()


# ═══════════════════════════════════════════
# SLIDE 12: BUSINESS IMPACT
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "10 — Business Impact of Undetected Tampering")
impacts = [
    ("🔓", "Data Breach", "Stolen credentials, customer data exfiltration", RED),
    ("💥", "Service Disruption", "Modified binaries causing system instability", AMBER),
    ("📋", "Compliance Failure", "Audit failures under PCI DSS, HIPAA, SOX", AMBER),
    ("📉", "Reputation Damage", "Loss of customer trust after breach disclosure", RED),
    ("💰", "Financial Loss", "Incident response costs, regulatory fines, lawsuits", RED),
    ("⏱️", "Operational Downtime", "Systems offline during investigation & recovery", AMBER),
]
for i, (icon, title, desc, color) in enumerate(impacts):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.3 + row * 2.8
    add_box(slide, x, y, 3.8, 2.3)
    add_text(slide, x, y + 0.2, 3.8, 0.5, icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, y + 0.8, 3.8, 0.4, title, font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.2, y + 1.3, 3.4, 0.7, desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 13: SYSTEM ARCHITECTURE DIAGRAM
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "11 — System Architecture Diagram")
# React panel
add_box(slide, 0.5, 1.3, 5.5, 5.5, fill_color=RGBColor(0xF0, 0xF8, 0xFF), border_color=CYAN)
add_text(slide, 0.5, 1.4, 5.5, 0.5, "React Frontend (:5173)", font_size=18, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
for i, comp in enumerate(["Dashboard + 5-Level Status", "File Table + Hash Viewer", "Alert Feed + Scan History", "File Browser Modal Popup", "Auto-Scan Scheduler"]):
    add_box(slide, 0.9, 2.1 + i * 0.85, 4.7, 0.65)
    add_text(slide, 1.1, 2.15 + i * 0.85, 4.3, 0.5, comp, font_size=13, color=WHITE)
# Arrow
add_text(slide, 6.1, 3.5, 1, 0.6, "⇄", font_size=36, color=CYAN, alignment=PP_ALIGN.CENTER)
add_text(slide, 6.1, 4.1, 1, 0.4, "REST API", font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)
# Python panel
add_box(slide, 7.3, 1.3, 5.5, 5.5, fill_color=RGBColor(0xF5, 0xF0, 0xFF), border_color=PURPLE)
add_text(slide, 7.3, 1.4, 5.5, 0.5, "Python Backend (:5000)", font_size=18, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
for i, comp in enumerate(["hashlib SHA-256 / MD5", "os.stat() File Metadata", "Real OS File Scanner", "Baseline Comparison Engine", "JSON Data Persistence"]):
    add_box(slide, 7.7, 2.1 + i * 0.85, 4.7, 0.65)
    add_text(slide, 7.9, 2.15 + i * 0.85, 4.3, 0.5, comp, font_size=13, color=PURPLE)


# ═══════════════════════════════════════════
# SLIDE 8: DATA FLOW DIAGRAM
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "12 — Data Flow Diagram")
flow_steps = [
    ("User Clicks Scan", CYAN, 4.5, 1.2),
    ("React → POST /api/scan", WHITE, 4, 2.2),
    ("Python Reads OS Files", PURPLE, 1.5, 3.4),
    ("Compute SHA-256 Hash", GREEN, 5.5, 3.4),
    ("Compare vs Baseline", AMBER, 4, 4.6),
    ("Match ✓  →  INTACT", GREEN, 1.5, 5.8),
    ("Mismatch ⚠  →  ALERT", RED, 7, 5.8),
]
for label, color, x, y in flow_steps:
    add_box(slide, x, y, 4, 0.7, border_color=color)
    add_text(slide, x, y + 0.05, 4, 0.6, label, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
# Arrows
for i in range(4):
    add_text(slide, 6.2, 1.85 + i * 1.15, 0.5, 0.4, "↓", font_size=20, color=DIM, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDES 9-15: DEMO SCREENSHOTS
# ═══════════════════════════════════════════
screenshots = [
    "Dashboard — System Overview + Status Indicator",
    "Dashboard — Auto-Scan Scheduler + Scan Summary",
    "Files Tab — Monitored Files Table",
    "Files Tab — Hash Comparison Detail",
    "Custom Scan — File Browser Modal Popup",
    "Alerts Tab — Security Alert Feed",
    "Report Tab — Auto-Generated Report",
]
for i, title in enumerate(screenshots):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, f"{16 + i} — Demo: {title}")
    add_screenshot_placeholder(slide, 0.8, 1.3, 11.5, 5.5, title)


# ═══════════════════════════════════════════
# SLIDE: RECOMMENDATIONS
# ═══════════════════════════════════════════
for section, items, color in [
    ("23 — Recommendations: Preventive", [
        "🔒  Use SELinux/AppArmor for mandatory access controls",
        "⚡  Enable real-time monitoring with inotify / auditd",
        "🏗️  Deploy read-only file systems for production",
        "💾  Store baselines on air-gapped storage",
        "🔐  Use SHA-256 exclusively — MD5/SHA-1 are broken",
    ], GREEN),
    ("24 — Recommendations: Detective", [
        "⏰  Schedule automated scans every 4-6 hours",
        "📊  Integrate FIM alerts with SIEM (Splunk, Elastic)",
        "🔍  Monitor permission changes (644 → 777 = compromise)",
        "🔗  Cross-reference FIM alerts with auth logs",
    ], CYAN),
    ("25 — Recommendations: Response", [
        "1.  Isolate affected systems immediately",
        "2.  Restore from trusted backups",
        "3.  Conduct forensic log analysis",
        "4.  Block C2 communication via firewall",
        "5.  Reset all exposed credentials",
        "6.  Document incident timeline",
        "7.  Update baseline after clean state",
    ], RED),
]:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, section)
    add_bullet_list(slide, 1, 1.3, 11, items, color=WHITE, font_size=16)


# ═══════════════════════════════════════════
# SLIDE: POC
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "26 — Proof of Concept")
poc_steps = [
    ("1", "Initialize Baseline", "Python hashes 20 system files → stores baseline.json", CYAN),
    ("2", "Run Integrity Scan", "Recalculates hashes → compares with baseline", GREEN),
    ("3", "Detect Violations", "Hash mismatch → DANGER alert with risk level", RED),
    ("4", "Inspect Details", "Click file → side-by-side SHA-256 comparison", PURPLE),
    ("5", "Auto-Scan", "Enable scheduler → system auto-scans at intervals", AMBER),
]
for i, (num, title, desc, color) in enumerate(poc_steps):
    add_text(slide, 1, 1.3 + i * 0.9, 0.5, 0.5, num, font_size=20, color=color, bold=True)
    add_text(slide, 1.7, 1.3 + i * 0.9, 4, 0.4, title, font_size=16, color=CYAN, bold=True)
    add_text(slide, 1.7, 1.7 + i * 0.9, 10, 0.4, desc, font_size=13, color=GRAY)
add_screenshot_placeholder(slide, 0.8, 5.8, 11.5, 1.2, "Workflow demo screenshot")


# ═══════════════════════════════════════════
# SLIDE: TOOLS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "27 — Tools & Technologies")
tools = [("Python 3.10+", "Backend server + file hashing"), ("hashlib", "SHA-256 / MD5 computation"),
         ("Flask", "REST API framework"), ("Flask-CORS", "Cross-origin access"),
         ("React.js (Vite)", "Frontend dashboard"), ("JavaScript ES6+", "Frontend logic"),
         ("JSON", "Data persistence"), ("Inter + JetBrains Mono", "Typography"),
         ("Git + GitHub", "Version control"), ("VS Code", "Code editor")]
for i, (tool, purpose) in enumerate(tools):
    y = 1.3 + i * 0.52
    add_text(slide, 1, y, 4, 0.4, tool, font_size=14, color=CYAN, bold=True)
    add_text(slide, 5, y, 7, 0.4, purpose, font_size=14, color=GRAY)
    if i < 9:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(y + 0.45), Inches(11), Inches(0.01))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BORDER
        shape.line.fill.background()


# ═══════════════════════════════════════════
# SLIDE: CHALLENGES & LEARNINGS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "28 — Challenges & Learnings")
# Challenges panel
add_box(slide, 0.8, 1.3, 5.5, 5, fill_color=RGBColor(0xFF, 0xF0, 0xF0), border_color=RGBColor(0xCC, 0x22, 0x44))
add_text(slide, 1.2, 1.4, 4.8, 0.4, "CHALLENGES", font_size=12, color=RED, bold=True)
challenges_list = [
    "Permission denied on protected system files (SAM)",
    "CORS configuration between React (:5173) and Flask (:5000)",
    "Cross-platform file paths (Windows \\ vs Linux /)",
    "Large file hashing efficiency (8KB chunked reads)",
    "Modal overlay z-index and positioning",
    "Auto-scan timer cleanup on component unmount",
]
for i, item in enumerate(challenges_list):
    add_text(slide, 1.2, 1.9 + i * 0.7, 4.8, 0.5, f"▸  {item}", font_size=13, color=WHITE)
# Learnings panel
add_box(slide, 6.8, 1.3, 5.5, 5, fill_color=RGBColor(0xF0, 0xFF, 0xF0), border_color=RGBColor(0x00, 0x8A, 0x4E))
add_text(slide, 7.2, 1.4, 4.8, 0.4, "KEY LEARNINGS", font_size=12, color=GREEN, bold=True)
learnings_list = [
    "Cryptographic hashing is essential for data integrity",
    "Full-stack architecture provides power + usability",
    "FIM is a critical compliance requirement in enterprise",
    "SHA-256 is current standard over MD5/SHA-1",
    "File metadata (perms, size) matters as much as content",
    "Separation of concerns: CSS/JSX/Python = clean code",
]
for i, item in enumerate(learnings_list):
    add_text(slide, 7.2, 1.9 + i * 0.7, 4.8, 0.5, f"✓  {item}", font_size=13, color=WHITE)


# ═══════════════════════════════════════════
# SLIDE: FUTURE ENHANCEMENTS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "29 — Future Enhancements")
futures = [
    ("Email / Slack Alerts", "Notify security teams on violations", "Planned", CYAN),
    ("Database Backend", "SQLite / PostgreSQL for better scale", "Planned", CYAN),
    ("Multi-host Monitoring", "Agent-based architecture for multiple servers", "Future", PURPLE),
    ("SIEM Integration", "CEF / Syslog export for Splunk, Elastic", "Future", PURPLE),
    ("CSV / PDF Export", "Download reports for compliance documentation", "Planned", CYAN),
    ("Windows Service", "Run as background service on startup", "Planned", CYAN),
    ("Real-time Watchdog", "inotify / ReadDirectoryChanges for instant alerts", "Future", PURPLE),
    ("File Diff Viewer", "Show exact content changes between versions", "Planned", CYAN),
]
for i, (title, desc, status, color) in enumerate(futures):
    y = 1.3 + i * 0.7
    add_text(slide, 0.8, y, 3, 0.4, title, font_size=14, color=CYAN, bold=True)
    add_text(slide, 4, y, 6, 0.4, desc, font_size=13, color=GRAY)
    add_text(slide, 10.5, y, 2, 0.4, status, font_size=11, color=color, bold=True)
    if i < 7:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y + 0.6), Inches(11.8), Inches(0.01))
        shape.fill.solid(); shape.fill.fore_color.rgb = BORDER; shape.line.fill.background()


# ═══════════════════════════════════════════
# SLIDE: API ENDPOINTS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "15 — REST API Endpoints")
api_data = [
    ("Method", "Endpoint", "Description"),
    ("GET", "/api/status", "System info + baseline status"),
    ("GET", "/api/files", "List monitored files with metadata"),
    ("POST", "/api/baseline/create", "Create trusted hash baseline"),
    ("POST", "/api/scan", "Run integrity scan vs baseline"),
    ("GET", "/api/alerts", "Get all security alerts"),
    ("POST", "/api/alerts/clear", "Clear alert history"),
    ("POST", "/api/browse", "Browse directory contents"),
    ("POST", "/api/custom-files", "Hash individual file"),
    ("POST", "/api/scan-folder", "Scan folder (recursive option)"),
    ("GET", "/api/report", "Generate security report"),
    ("GET", "/api/history", "Get all scan history"),
]
for i, row in enumerate(api_data):
    y = 1.3 + i * 0.48
    is_header = i == 0
    c0 = CYAN if is_header else (GREEN if row[0] == "POST" else CYAN)
    add_text(slide, 0.8, y, 1.5, 0.35, row[0], font_size=12 if not is_header else 11, color=c0, bold=True)
    add_text(slide, 2.5, y, 4, 0.35, row[1], font_size=12, color=CYAN if is_header else WHITE, bold=is_header, font_name="Consolas" if not is_header else "Segoe UI")
    add_text(slide, 6.8, y, 5, 0.35, row[2], font_size=12, color=CYAN if is_header else GRAY, bold=is_header)
    if i > 0:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y + 0.42), Inches(11.8), Inches(0.01))
        shape.fill.solid(); shape.fill.fore_color.rgb = BORDER; shape.line.fill.background()


# ═══════════════════════════════════════════
# SLIDE: REFERENCES
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_slide_title(slide, "30 — References")
refs = [
    "[1]  OSSEC Syscheck Documentation — ossec.net",
    "[2]  AccuKnox — Top FIM Tools 2025 — accuknox.com",
    "[3]  Comparitech — Best FIM Tools — comparitech.com",
    "[4]  SolarWinds — FIM Software — solarwinds.com",
    "[5]  WafaTech — Mastering AIDE — wafatech.sa",
    "[6]  OSSEC GitHub — Hash Discussion — github.com/ossec",
    "[7]  MojoAuth — MD5 vs SHA-256 — mojoauth.com",
    "[8]  DNSstuff — FIM for Enterprises — dnsstuff.com",
    "[9]  TrustRadius — FIM Reviews — trustradius.com",
    "[10] PCI DSS v4.0 Requirement 11.5 — pcisecuritystandards.org",
]
for i, ref in enumerate(refs):
    add_text(slide, 1, 1.3 + i * 0.5, 11, 0.4, ref, font_size=13, color=GRAY)


# ═══════════════════════════════════════════
# SLIDE: THANK YOU
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0, 1.5, 13.3, 0.8, "🛡️", font_size=60, alignment=PP_ALIGN.CENTER)
add_text(slide, 0, 2.5, 13.3, 0.8, "Thank You!", font_size=44, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 0, 3.4, 13.3, 0.5, "File Integrity Guardian v1.0", font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, 0, 4.2, 13.3, 0.4, "Full-Stack FIM Tool  •  Python + React  •  SHA-256 / MD5", font_size=14, color=DIM, alignment=PP_ALIGN.CENTER)
add_text(slide, 0, 4.8, 13.3, 0.4, "Built by Gaurav Kumar", font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
add_box(slide, 4, 5.5, 5.3, 1, fill_color=RGBColor(0xF0, 0xF8, 0xFF), border_color=CYAN)
add_text(slide, 4, 5.6, 5.3, 0.4, "Questions & Discussion", font_size=16, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 4, 6.05, 5.3, 0.3, "github.com/YOUR_USERNAME/file-integrity-guardian", font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════
output = "File_Integrity_Guardian_Presentation.pptx"
prs.save(output)
print(f"""
╔═══════════════════════════════════════════════════╗
║  ✅ PowerPoint Generated Successfully!             ║
╠═══════════════════════════════════════════════════╣
║  File:   {output:<40}║
║  Slides: {len(prs.slides):<40}║
║  Path:   {os.path.abspath(output)[:40]:<40}║
╚═══════════════════════════════════════════════════╝

  Now open it in PowerPoint and paste your screenshots
  into the labeled placeholder slides!
""")